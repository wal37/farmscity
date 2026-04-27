from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/Users/wale/Desktop/farmer/farmerweb/FARMSCITY_PITCH_DECK.pptx")


def rgb(value: str) -> RGBColor:
    value = value.replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


NIGHT = rgb("#111511")
OLIVE = rgb("#3C5E31")
LEAF = rgb("#6E9633")
STRAW = rgb("#D7B15D")
TERRA = rgb("#B96B3C")
MIST = rgb("#E4E8DD")
PAPER = rgb("#F7F4EC")
WHITE = rgb("#FFFFFF")
INK = rgb("#1D211B")
MUTED = rgb("#62685B")
LINE = rgb("#D7DCCE")
PALE = rgb("#EFE8D8")
SAGE = rgb("#D5E4C3")
BLUSH = rgb("#E8D4C5")
SKY = rgb("#B8D1D7")
SAND = rgb("#E6D8BE")


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_round(slide, left, top, width, height, fill, line=None, radius=0.04):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.adjustments[0] = radius
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=20,
    color=INK,
    bold=False,
    font="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_label(slide, left, top, text, fill=INK, color=WHITE):
    width = max(1.5, 0.102 * len(text) + 0.45)
    add_round(slide, left, top, Inches(width), Inches(0.38), fill, None, 0.14)
    add_text(slide, left + Inches(0.16), top + Inches(0.07), Inches(width - 0.18), Inches(0.2), text.upper(), size=10, color=color, bold=True)


def add_card(slide, left, top, width, height, title, body, fill, title_color=INK, body_color=MUTED):
    add_round(slide, left, top, width, height, fill, None, 0.04)
    add_text(slide, left + Inches(0.22), top + Inches(0.22), width - Inches(0.38), Inches(0.34), title, size=18, color=title_color, bold=True)
    add_text(slide, left + Inches(0.22), top + Inches(0.72), width - Inches(0.4), height - Inches(0.92), body, size=12, color=body_color)


def add_logo(slide, left, top, size=0.7):
    scale = size / 0.82
    add_round(slide, left, top, Inches(0.82 * scale), Inches(0.82 * scale), rgb("#132A1A"), None, 0.18)
    add_round(slide, left + Inches(0.14 * scale), top + Inches(0.5 * scale), Inches(0.54 * scale), Inches(0.12 * scale), rgb("#D88A34"), None, 0.08)
    add_round(slide, left + Inches(0.18 * scale), top + Inches(0.56 * scale), Inches(0.46 * scale), Inches(0.07 * scale), rgb("#F5E9C9"), None, 0.08)
    add_rect(slide, left + Inches(0.18 * scale), top + Inches(0.31 * scale), Inches(0.18 * scale), Inches(0.16 * scale), rgb("#2F7F52"))
    add_rect(slide, left + Inches(0.36 * scale), top + Inches(0.28 * scale), Inches(0.14 * scale), Inches(0.19 * scale), rgb("#88C062"))
    add_rect(slide, left + Inches(0.5 * scale), top + Inches(0.31 * scale), Inches(0.14 * scale), Inches(0.16 * scale), rgb("#4E7BB9"))
    add_rect(slide, left + Inches(0.28 * scale), top + Inches(0.14 * scale), Inches(0.08 * scale), Inches(0.11 * scale), rgb("#88C062"))
    add_rect(slide, left + Inches(0.39 * scale), top + Inches(0.11 * scale), Inches(0.08 * scale), Inches(0.12 * scale), rgb("#F0C24B"))
    add_rect(slide, left + Inches(0.5 * scale), top + Inches(0.15 * scale), Inches(0.08 * scale), Inches(0.1 * scale), rgb("#7EB7C7"))


def add_page(slide, num, fill=PAPER):
    add_rect(slide, 0, 0, Inches(13.333), Inches(7.5), fill)
    add_rect(slide, Inches(0.78), Inches(0.72), Inches(11.78), Inches(0.03), INK if fill != NIGHT else WHITE)
    add_rect(slide, Inches(0.78), Inches(6.82), Inches(11.78), Inches(0.03), INK if fill != NIGHT else WHITE)
    add_text(slide, Inches(12.02), Inches(0.26), Inches(0.38), Inches(0.18), f"{num:02d}", size=10, color=INK if fill != NIGHT else WHITE, bold=True, align=PP_ALIGN.RIGHT)


def cover(slide, num):
    add_page(slide, num, fill=NIGHT)
    add_rect(slide, Inches(7.72), Inches(0.75), Inches(4.84), Inches(6.07), OLIVE)
    add_rect(slide, Inches(8.25), Inches(1.22), Inches(3.78), Inches(0.92), STRAW)
    add_rect(slide, Inches(8.65), Inches(2.46), Inches(2.95), Inches(0.78), TERRA)
    add_rect(slide, Inches(8.25), Inches(3.58), Inches(3.78), Inches(2.66), LEAF)
    add_logo(slide, Inches(1.02), Inches(1.02), 0.86)
    add_label(slide, Inches(1.02), Inches(2.08), "Agricultural operating network", fill=WHITE, color=INK)
    add_text(slide, Inches(1.02), Inches(2.68), Inches(5.92), Inches(0.82), "FarmsCity", size=36, color=WHITE, bold=True)
    add_text(slide, Inches(1.02), Inches(3.66), Inches(6.0), Inches(1.48), "A networked platform for crop planning, livestock operations, weather response, field records, and performance across agricultural teams.", size=21, color=MIST)
    add_text(slide, Inches(1.02), Inches(6.18), Inches(2.8), Inches(0.22), "Pitch deck | April 2026", size=12, color=WHITE, bold=True)


def problem(slide, num):
    add_page(slide, num)
    add_label(slide, Inches(1.0), Inches(1.04), "Problem", fill=TERRA)
    add_text(slide, Inches(1.0), Inches(1.58), Inches(11.0), Inches(0.82), "Agricultural operations still fail when timing, records, and coordination are separated.", size=29, bold=True)
    add_rect(slide, Inches(1.0), Inches(2.66), Inches(11.1), Inches(0.06), INK)
    add_card(slide, Inches(1.0), Inches(3.18), Inches(3.45), Inches(2.18), "Fragmented records", "Crop, livestock, labor, and weather activity are often tracked in separate places, which breaks operational continuity.", MIST)
    add_card(slide, Inches(4.82), Inches(3.18), Inches(3.45), Inches(2.18), "Slow visibility", "Managers and field teams rarely share the same real-time operational picture across sites.", BLUSH)
    add_card(slide, Inches(8.64), Inches(3.18), Inches(3.45), Inches(2.18), "Weak network learning", "Performance patterns disappear when farms and programs cannot compare clean histories over time.", PAPER)


def why_now(slide, num):
    add_page(slide, num, fill=MIST)
    add_label(slide, Inches(1.0), Inches(1.02), "Why now", fill=OLIVE)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.4), Inches(0.82), "The next agricultural edge is operational intelligence, not just more raw data.", size=28, bold=True)
    add_text(slide, Inches(1.0), Inches(2.48), Inches(6.2), Inches(0.96), "As cooperatives, advisory teams, and multi-site operators grow, the bottleneck shifts toward coordination, follow-through, and shared historical memory.", size=15, color=MUTED)
    add_rect(slide, Inches(7.7), Inches(1.0), Inches(4.28), Inches(5.2), WHITE)
    add_card(slide, Inches(8.02), Inches(1.34), Inches(3.62), Inches(1.0), "Shift 1", "Weather and input volatility make timing decisions more expensive.", STRAW)
    add_card(slide, Inches(8.02), Inches(2.84), Inches(3.62), Inches(1.0), "Shift 2", "Agricultural programs now need cleaner reporting from field to network level.", LEAF, title_color=WHITE, body_color=WHITE)
    add_card(slide, Inches(8.02), Inches(4.34), Inches(3.62), Inches(1.0), "Shift 3", "Farm networks need shared visibility instead of isolated farm-by-farm tools.", TERRA, title_color=WHITE, body_color=WHITE)


def platform(slide, num):
    add_page(slide, num)
    add_rect(slide, Inches(0.82), Inches(1.0), Inches(11.7), Inches(1.22), OLIVE)
    add_label(slide, Inches(1.08), Inches(1.26), "Platform", fill=WHITE, color=INK)
    add_text(slide, Inches(1.08), Inches(2.6), Inches(10.6), Inches(0.82), "FarmsCity is a coordinated operating layer, not a single-purpose farm logbook.", size=28, bold=True)
    add_text(slide, Inches(1.08), Inches(3.48), Inches(10.2), Inches(0.9), "The platform combines crop planning, livestock records, weather-linked timing, field activity, and performance review into one working system for farms and agricultural networks.", size=15, color=MUTED)
    add_card(slide, Inches(1.08), Inches(4.58), Inches(2.55), Inches(1.72), "Crop layer", "Rotation maps, seasonal calendars, treatments, and harvest tracking.", WHITE)
    add_card(slide, Inches(3.9), Inches(4.58), Inches(2.55), Inches(1.72), "Livestock layer", "Health, breeding, feed, and movement records.", PALE, title_color=INK, body_color=MUTED)
    add_card(slide, Inches(6.72), Inches(4.58), Inches(2.55), Inches(1.72), "Weather layer", "Rainfall alerts, field timing, and task shifts.", WHITE)
    add_card(slide, Inches(9.54), Inches(4.58), Inches(2.55), Inches(1.72), "Records layer", "Inputs, labor, audit history, and financial memory.", PALE)


def modules(slide, num):
    add_page(slide, num, fill=PAPER)
    add_label(slide, Inches(1.0), Inches(1.04), "Core modules", fill=LEAF)
    add_text(slide, Inches(1.0), Inches(1.58), Inches(6.9), Inches(0.82), "The product is organized around how agricultural work actually happens.", size=28, bold=True)
    add_rect(slide, Inches(1.0), Inches(2.8), Inches(11.0), Inches(0.06), OLIVE)
    add_card(slide, Inches(1.0), Inches(3.28), Inches(3.45), Inches(2.16), "Crop planning", "Field maps, crop rotation, planting windows, and seasonal task sequencing.", MIST)
    add_card(slide, Inches(4.82), Inches(3.28), Inches(3.45), Inches(2.16), "Livestock operations", "Animal history, breeding cycles, feed patterns, and health interventions.", BLUSH)
    add_card(slide, Inches(8.64), Inches(3.28), Inches(3.45), Inches(2.16), "Records and performance", "Inputs, field actions, losses, yields, and cross-season comparison.", WHITE)


def workflow(slide, num):
    add_page(slide, num, fill=SAGE)
    add_label(slide, Inches(1.0), Inches(1.02), "Workflow", fill=INK)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.3), Inches(0.82), "The operating rhythm is plan, run, respond, and learn.", size=28, bold=True)
    add_rect(slide, Inches(1.0), Inches(4.12), Inches(10.95), Inches(0.08), INK)
    steps = [
        (1.16, "01", "Plan", "Set fields, herds, targets, and weather-sensitive schedules.", STRAW),
        (4.06, "02", "Run", "Record labor, inputs, animal events, and field conditions.", TERRA),
        (6.96, "03", "Respond", "Adjust quickly around rainfall, disease pressure, and supply changes.", OLIVE),
        (9.86, "04", "Learn", "Compare fields, breeds, seasons, and interventions across the network.", LEAF),
    ]
    for x, no, title, body, color in steps:
        add_rect(slide, Inches(x), Inches(3.88), Inches(0.18), Inches(0.54), color)
        add_text(slide, Inches(x + 0.28), Inches(3.32), Inches(2.0), Inches(0.22), no, size=11, color=color, bold=True)
        add_text(slide, Inches(x + 0.28), Inches(4.56), Inches(2.0), Inches(0.28), title, size=20, bold=True)
        add_text(slide, Inches(x + 0.28), Inches(5.0), Inches(2.02), Inches(0.86), body, size=12, color=MUTED)


def records(slide, num):
    add_page(slide, num)
    add_label(slide, Inches(1.0), Inches(1.02), "Records", fill=STRAW, color=INK)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.5), Inches(0.82), "Every action should feed a record system the wider network can use.", size=28, bold=True)
    add_text(slide, Inches(1.0), Inches(2.48), Inches(6.3), Inches(0.96), "FarmsCity preserves field history, herd events, treatments, harvests, and operator activity so decisions are made against real history instead of memory.", size=15, color=MUTED)
    add_rect(slide, Inches(7.86), Inches(0.98), Inches(4.22), Inches(5.62), OLIVE)
    add_card(slide, Inches(8.16), Inches(1.28), Inches(3.62), Inches(1.04), "Field history", "Planting, spraying, harvest, and labor remain linked.", WHITE)
    add_card(slide, Inches(8.16), Inches(2.76), Inches(3.62), Inches(1.04), "Livestock history", "Health, breeding, and feed patterns stay reviewable.", STRAW)
    add_card(slide, Inches(8.16), Inches(4.24), Inches(3.62), Inches(1.04), "Network reporting", "Managers and programs can compare clean activity records.", WHITE)


def customers(slide, num):
    add_page(slide, num, fill=MIST)
    add_label(slide, Inches(1.0), Inches(1.04), "Customers", fill=LEAF)
    add_text(slide, Inches(1.0), Inches(1.58), Inches(6.7), Inches(0.82), "FarmsCity serves both individual operators and larger agricultural networks.", size=28, bold=True)
    add_card(slide, Inches(1.0), Inches(3.18), Inches(2.74), Inches(1.88), "Farm owners", "Run crops, livestock, and records from one operating surface.", WHITE)
    add_card(slide, Inches(3.98), Inches(3.18), Inches(2.74), Inches(1.88), "Cooperatives", "Standardize field visibility and reporting across members.", PAPER)
    add_card(slide, Inches(6.96), Inches(3.18), Inches(2.74), Inches(1.88), "Extension teams", "Track follow-up, support activity, and field histories.", WHITE)
    add_card(slide, Inches(9.94), Inches(3.18), Inches(2.14), Inches(1.88), "Programs", "Support deployment, oversight, and performance review.", PAPER)


def moat(slide, num):
    add_page(slide, num)
    add_label(slide, Inches(1.0), Inches(1.02), "Differentiation", fill=TERRA)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.8), Inches(0.82), "The moat is operational coordination across farms, not isolated logging.", size=28, bold=True)
    add_rect(slide, Inches(1.0), Inches(2.72), Inches(11.1), Inches(0.06), INK)
    add_card(slide, Inches(1.0), Inches(3.22), Inches(3.45), Inches(2.04), "Network-aware design", "Built for single farms and multi-site operating structures.", MIST)
    add_card(slide, Inches(4.82), Inches(3.22), Inches(3.45), Inches(2.04), "Decision timing", "Weather, field, and livestock context stay close to the action.", BLUSH)
    add_card(slide, Inches(8.64), Inches(3.22), Inches(3.45), Inches(2.04), "Persistent memory", "Historical records stay usable across seasons and staff changes.", PAPER)


def market(slide, num):
    add_page(slide, num, fill=NIGHT)
    add_rect(slide, Inches(0.78), Inches(0.72), Inches(11.78), Inches(0.03), WHITE)
    add_rect(slide, Inches(0.78), Inches(6.82), Inches(11.78), Inches(0.03), WHITE)
    add_label(slide, Inches(1.0), Inches(1.02), "Market", fill=WHITE, color=INK)
    add_text(slide, Inches(1.0), Inches(1.56), Inches(6.8), Inches(0.82), "The market spans operators, cooperatives, advisory networks, and agricultural programs.", size=28, color=WHITE, bold=True)
    add_text(slide, Inches(1.0), Inches(2.5), Inches(6.6), Inches(0.94), "The same platform can start as a farm operating system and grow into a coordination layer for larger agricultural ecosystems.", size=15, color=MIST)
    add_card(slide, Inches(7.72), Inches(1.34), Inches(4.0), Inches(1.1), "Direct farm subscriptions", "Operational software for owners and managers.", STRAW)
    add_card(slide, Inches(7.72), Inches(2.86), Inches(4.0), Inches(1.1), "Cooperative deployments", "Shared records and reporting across member farms.", WHITE)
    add_card(slide, Inches(7.72), Inches(4.38), Inches(4.0), Inches(1.1), "Program and advisory rollouts", "Visibility and support for field-level operations.", SAGE)


def business(slide, num):
    add_page(slide, num)
    add_label(slide, Inches(1.0), Inches(1.04), "Business model", fill=LEAF)
    add_text(slide, Inches(1.0), Inches(1.58), Inches(6.1), Inches(0.82), "Software revenue, deployment support, and network operations tooling.", size=28, bold=True)
    add_card(slide, Inches(1.0), Inches(3.28), Inches(3.5), Inches(1.94), "Farm subscriptions", "Charge for access by farm, manager, or operating team.", WHITE)
    add_card(slide, Inches(4.82), Inches(3.28), Inches(3.5), Inches(1.94), "Cooperative licenses", "Support group reporting, monitoring, and deployment workflows.", BLUSH)
    add_card(slide, Inches(8.64), Inches(3.28), Inches(3.5), Inches(1.94), "Implementation support", "Rollout, workflow design, and reporting support for larger operators.", MIST)


def roadmap(slide, num):
    add_page(slide, num, fill=SAND)
    add_label(slide, Inches(1.0), Inches(1.04), "Roadmap", fill=OLIVE)
    add_text(slide, Inches(1.0), Inches(1.58), Inches(6.6), Inches(0.82), "The build path starts with operations and expands into network intelligence.", size=28, bold=True)
    add_rect(slide, Inches(1.0), Inches(4.02), Inches(11.05), Inches(0.08), INK)
    phases = [
        (1.12, "Phase 1", "Core crop, livestock, weather, and record modules for direct farm use.", STRAW),
        (4.72, "Phase 2", "Inventory workflows, stronger alerts, and cross-site comparison.", TERRA),
        (8.32, "Phase 3", "Deeper network reporting, advisory coordination, and predictive planning.", LEAF),
    ]
    for x, title, body, color in phases:
        add_rect(slide, Inches(x), Inches(3.8), Inches(0.16), Inches(0.56), color)
        add_text(slide, Inches(x + 0.26), Inches(2.96), Inches(2.8), Inches(0.28), title, size=20, bold=True)
        add_text(slide, Inches(x + 0.26), Inches(4.42), Inches(2.62), Inches(0.86), body, size=12, color=MUTED)


def closing(slide, num):
    add_page(slide, num, fill=WHITE)
    add_rect(slide, Inches(0.86), Inches(0.96), Inches(11.58), Inches(5.9), LEAF)
    add_logo(slide, Inches(1.22), Inches(1.36), 0.8)
    add_label(slide, Inches(1.22), Inches(2.34), "Closing", fill=WHITE, color=INK)
    add_text(slide, Inches(1.22), Inches(2.9), Inches(6.4), Inches(0.82), "FarmsCity gives agriculture a cleaner operating network.", size=31, color=WHITE, bold=True)
    add_text(slide, Inches(1.22), Inches(3.88), Inches(6.3), Inches(1.0), "The product connects crops, livestock, weather, records, and performance so farms and agricultural networks can act earlier and learn faster.", size=18, color=MIST)
    add_rect(slide, Inches(8.28), Inches(1.42), Inches(3.42), Inches(1.0), STRAW)
    add_rect(slide, Inches(8.88), Inches(2.74), Inches(2.82), Inches(0.82), TERRA)
    add_rect(slide, Inches(8.28), Inches(3.94), Inches(3.42), Inches(1.78), WHITE)
    add_text(slide, Inches(8.54), Inches(4.36), Inches(2.8), Inches(0.32), "FarmsCity", size=22, color=INK, bold=True)
    add_text(slide, Inches(8.54), Inches(4.84), Inches(2.76), Inches(0.66), "Agricultural operations platform for crops, livestock, records, and weather-aware network decisions.", size=12, color=MUTED)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    slides = [
        cover,
        problem,
        why_now,
        platform,
        modules,
        workflow,
        records,
        customers,
        moat,
        market,
        business,
        roadmap,
        closing,
    ]

    for idx, fn in enumerate(slides, start=1):
        slide = prs.slides.add_slide(layout)
        fn(slide, idx)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
